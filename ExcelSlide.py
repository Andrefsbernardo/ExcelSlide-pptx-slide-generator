import pandas as pd
from pptx import Presentation
from pptx.util import Pt

# Load the spreadsheet
spreadsheet_path = './RationaleFile.xlsx'
data = pd.read_excel(spreadsheet_path)

# Load the PowerPoint template
template_path = './TemplateFile.pptx'
presentation = Presentation(template_path)

# Function to replace placeholders while keeping the format consistent
def replace_placeholders(slide, placeholders):
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            for placeholder, value in placeholders.items():
                if placeholder in shape.text:
                    original_text = shape.text
                    shape.text = shape.text.replace(placeholder, str(value))
                    
                    # Retain the font size and font style
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if placeholder in original_text:
                                    # Keep the original font size and font style
                                    if run.font.size:
                                        run.font.size = Pt(run.font.size.pt)
                                    if run.font.name:
                                        run.font.name = run.font.name

# Function to duplicate a slide
def duplicate_slide(presentation, slide_index):
    slide = presentation.slides[slide_index]
    slide_layout = slide.slide_layout
    new_slide = presentation.slides.add_slide(slide_layout)

    # Copy each shape from the existing slide to the new one
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.has_text_frame:
            # Create a new textbox in the same position with the same text
            new_shape = new_slide.shapes.add_textbox(
                left=shape.left, top=shape.top, width=shape.width, height=shape.height
            )
            new_text_frame = new_shape.text_frame
            for paragraph in shape.text_frame.paragraphs:
                new_paragraph = new_text_frame.add_paragraph()
                new_paragraph.text = paragraph.text

                # Retain font style and size
                if paragraph.runs:
                    for run, new_run in zip(paragraph.runs, new_paragraph.runs):
                        new_run.font.size = run.font.size
                        new_run.font.name = run.font.name

    return new_slide

# Iterate over each row in the spreadsheet
for index, row in data.iterrows():
    # Duplicate the first slide (assuming it's the template slide) to create a new one
    new_slide = duplicate_slide(presentation, 0)
    
    # Map data from the current row to the placeholders in the new slide
    placeholders = {
        '[First Name]': row['First Name'],
        '[Lookup ID]': str(row['Lookup ID']),
        '[Last Name]': row['Last Name'],
        '[Title]': row['Title'],
        '[Spouse Name]': row['Spouse Name'],
        '[SK Patient Connection (Y/N)]': row['SK Patient Connection (Y/N)'],
        '[Other SK affiliations/ board connections]': row['Other SK affiliations/ board connections'],
        '[Giving to SKF]': row['Giving to SKF'],
        '[Most recent SK Major Gift (gifts at/over $10K)]': row['Most recent SK Major Gift (gifts at/over $10K)'],
        '[Suspected areas of interest]': row['Suspected areas of interest'],
        '[Philanthropy examples - $5M+]': row['Philanthropy examples - $5M+'],
        '[Known Influencers]': row['Known Influencers'],
        '[Recent Engagement]': row['Recent Engagement'],
        '[Research Comments]': row['Research Comments']
    }
    
    # Replace placeholders in the new slide
    replace_placeholders(new_slide, placeholders)

# Save the final presentation
output_path = './output_presentation.pptx'
presentation.save(output_path)

print(f"Presentation saved to {output_path}")